# -*- coding: utf-8 -*-
"""GOOYAのAI活用図鑑 PowerPoint デザイン部品（python-pptx 用）。

サイト（gooya-ai-atlas/src/index.html）の CSS トークンと1対1対応。
使い方は references/design-spec.md を参照。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

COLORS = {
    "BRAND": RGBColor(0x6B, 0xBF, 0x3F),        # --c-brand
    "BRAND_DARK": RGBColor(0x55, 0xA0, 0x2F),   # --c-brand-dark
    "BRAND_SOFT": RGBColor(0xEA, 0xF7, 0xE1),   # --c-brand-soft
    "BRAND_TINT": RGBColor(0xF4, 0xFB, 0xEF),   # --c-brand-tint
    "TEXT": RGBColor(0x1E, 0x2A, 0x22),         # --c-text
    "TEXT_SUB": RGBColor(0x5E, 0x6E, 0x64),     # --c-text-sub
    "TEXT_MUTE": RGBColor(0x8A, 0x99, 0x90),    # --c-text-mute
    "BORDER": RGBColor(0xE3, 0xE9, 0xE3),       # --c-border
    "SURFACE_ALT": RGBColor(0xF7, 0xFA, 0xF7),  # --c-surface-alt
    "DANGER": RGBColor(0xC2, 0x45, 0x3B),       # --c-danger
    "WHITE": RGBColor(0xFF, 0xFF, 0xFF),
}

FONT = "Yu Gothic"  # 游ゴシック。set_text() が Latin/EastAsian 両方に設定する

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs, bg=None):
    """白紙スライドを追加。bg に COLORS のキー名か RGBColor で背景色を指定できる。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        color = COLORS[bg] if isinstance(bg, str) else bg
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color
    return slide


def _apply_font(run, size, bold, color, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = COLORS[color] if isinstance(color, str) else color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", FONT)


def set_text(tf_or_para, text, size=14, bold=False, color="TEXT",
             align=PP_ALIGN.LEFT, line_spacing=1.15, space_after=0):
    """段落にテキストを設定。text は str か、(文字列, 上書きスタイルdict) のリスト。

    リストを渡すと1段落内で部分ごとにスタイルを変えられる:
      set_text(p, [("連携方式  ", {"bold": True}), ("直接参照", {"color": "TEXT_SUB"})])
    """
    para = tf_or_para.paragraphs[0] if hasattr(tf_or_para, "paragraphs") else tf_or_para
    para.alignment = align
    if line_spacing:
        para.line_spacing = line_spacing
    para.space_after = Pt(space_after)
    parts = [(text, {})] if isinstance(text, str) else text
    for t, ov in parts:
        run = para.add_run()
        run.text = t
        _apply_font(run, ov.get("size", size), ov.get("bold", bold),
                    ov.get("color", color), ov.get("italic", False))
    return para


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    """空のテキストボックス。内部余白ゼロ・折り返しあり。set_text と組み合わせて使う。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return box


def add_para(tf, **kwargs):
    """テキストフレームに段落を追加して set_text を適用（2段落目以降用）。"""
    para = tf.add_paragraph()
    return set_text(para, **kwargs)


def _bullet(para, char="•", color="BRAND", indent_in=0.16):
    """段落に色付き行頭記号を付ける。space_before/space_after は先に設定しておくこと。"""
    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", str(Inches(indent_in)))
    pPr.set("indent", str(-Inches(indent_in)))
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgb = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", "%02X%02X%02X" % tuple(COLORS[color] if isinstance(color, str) else color))
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


def add_bullets(tf, items, size=13, color="TEXT", space_after=6, line_spacing=1.15,
                bullet_color="BRAND"):
    """テキストフレームへ箇条書きを流し込む。items は str のリスト。

    最初の項目はフレームの先頭段落を使うので、空のテキストフレームに対して呼ぶこと。
    """
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_text(para, item, size=size, color=color,
                 line_spacing=line_spacing, space_after=space_after)
        _bullet(para, color=bullet_color)


def _round_rect(slide, x, y, w, h, radius_in):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = min(0.5, radius_in / min(w, h))
    shape.shadow.inherit = False
    return shape


def add_card(slide, x, y, w, h, fill="SURFACE_ALT", border="BORDER", radius_in=0.12):
    """角丸カード。戻り値のシェイプに直接テキストは入れず、上に add_textbox を重ねる。"""
    card = _round_rect(slide, x, y, w, h, radius_in)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = COLORS[border] if isinstance(border, str) else border
        card.line.width = Pt(1)
    return card


def add_pill(slide, x, y, w, h, text, size=10.5, fill="BRAND_SOFT",
             text_color="BRAND_DARK", border=None, bold=True):
    """ピルバッジ（完全角丸）。カテゴリ・タグ用。

    アウトライン版は fill="WHITE", border="BORDER", text_color="TEXT_SUB"。
    """
    pill = _round_rect(slide, x, y, w, h, radius_in=h)  # adjustment=0.5 になる
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    if border is None:
        pill.line.fill.background()
    else:
        pill.line.color.rgb = COLORS[border] if isinstance(border, str) else border
        pill.line.width = Pt(1)
    tf = pill.text_frame
    tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, text, size=size, bold=bold, color=text_color,
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    return pill


def add_number_badge(slide, x, y, d, number, size=14):
    """BRAND 塗りの正円 + 白太字の数字。d は直径（インチ）。"""
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS["BRAND"]
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, str(number), size=size, bold=True, color="WHITE",
             align=PP_ALIGN.CENTER, line_spacing=1.0)
    return circle


def add_icon_badge(slide, x, y, d, image_path, icon_ratio=0.52, fill="BRAND"):
    """BRAND 塗りの正円の上に白抜きアイコン画像を重ねる。d は円の直径（インチ）。

    白抜き（白一色）の PNG アイコンを前提にした部品。元資料から抜き出したアイコンを
    再利用するときは、円の色だけ図鑑のブランド緑に変えれば意匠が揃う。
    """
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS[fill] if isinstance(fill, str) else fill
    circle.line.fill.background()
    circle.shadow.inherit = False
    s = d * icon_ratio
    pic = slide.shapes.add_picture(
        image_path, Inches(x + (d - s) / 2), Inches(y + (d - s) / 2),
        Inches(s), Inches(s))
    return circle, pic


def add_chevron(slide, x, y, w=0.28, h=0.28):
    """フロー図の区切り用シェブロン（BRAND 色・枠なし）。"""
    ch = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    ch.fill.solid()
    ch.fill.fore_color.rgb = COLORS["BRAND"]
    ch.line.fill.background()
    ch.shadow.inherit = False
    return ch


def add_source(slide, text, x=0.5, y=7.08, w=9.0):
    """出典フッター。左下固定・9pt・TEXT_MUTE。"""
    box = add_textbox(slide, x, y, w, 0.3)
    set_text(box.text_frame, text, size=9, color="TEXT_MUTE", line_spacing=1.0)
    return box
